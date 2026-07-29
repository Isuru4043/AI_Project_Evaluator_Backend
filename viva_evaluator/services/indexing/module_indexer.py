"""
Module Materials Indexer — Parses lecture notes/slides and builds a FAISS index.
"""
import os
import json
import logging
from typing import List, Dict

logger = logging.getLogger(__name__)

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

def _extract_text_from_pdf(filepath: str) -> str:
    if not fitz:
        logger.error("PyMuPDF (fitz) is not installed.")
        return ""
    text = []
    try:
        with fitz.open(filepath) as doc:
            for page in doc:
                text.append(page.get_text())
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
    return "\n".join(text)

def _extract_text_from_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except ImportError:
        logger.error("python-pptx is not installed.")
        return ""
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        return ""

def index_module_material(material) -> bool:
    """
    Extracts text, creates chunks, and saves to storage.
    Updates the material processing status.
    """
    from core.models import ModuleMaterial
    from django.core.files.storage import default_storage
    from django.core.files.base import ContentFile
    import tempfile
    import requests
    import faiss
    import numpy as np
    from viva_evaluator.services.rag.chunking import chunk_text
    from viva_evaluator.services.rag.embeddings import embed_texts, EMBEDDING_DIM
    
    try:
        import urllib.parse
        file_url = material.file_url
        parsed = urllib.parse.urlparse(file_url)
        path = parsed.path
        if 'module_materials/' in path:
            path = path[path.index('module_materials/'):]
        else:
            path = file_url.replace('/media/', '')
            
        with default_storage.open(path, 'rb') as f:
            content = f.read()

        ext = os.path.splitext(material.original_filename)[1].lower()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
            temp_file.write(content)
            temp_path = temp_file.name

        text = ""
        if ext == '.pdf':
            text = _extract_text_from_pdf(temp_path)
        elif ext == '.pptx':
            text = _extract_text_from_pptx(temp_path)
        
        os.unlink(temp_path)

        if not text.strip():
            logger.warning(f"No text extracted from {material.original_filename}")
            material.processing_status = ModuleMaterial.ProcessingStatus.FAILED
            material.save()
            return False

        # Use native chunking
        chunks = chunk_text(text, source='module', section=material.original_filename)
        
        if not chunks:
            logger.warning("No chunks generated")
            material.processing_status = ModuleMaterial.ProcessingStatus.FAILED
            material.save()
            return False

        texts = [c['text'] for c in chunks]
        vectors = embed_texts(texts)
        
        # Build FAISS index
        index = faiss.IndexFlatIP(EMBEDDING_DIM)
        index.add(vectors.astype(np.float32))
        
        # Serialize FAISS index
        serialized = faiss.serialize_index(index)
        faiss_blob = bytes(serialized)
        
        # Save chunks and faiss blob to storage (no DB migration needed)
        chunks_json = json.dumps(chunks).encode('utf-8')
        
        chunks_path = f"module_materials/{material.id}_chunks.json"
        faiss_path = f"module_materials/{material.id}_faiss.bin"
        
        # Overwrite if exists
        if default_storage.exists(chunks_path):
            default_storage.delete(chunks_path)
        if default_storage.exists(faiss_path):
            default_storage.delete(faiss_path)
            
        default_storage.save(chunks_path, ContentFile(chunks_json))
        default_storage.save(faiss_path, ContentFile(faiss_blob))

        material.processing_status = ModuleMaterial.ProcessingStatus.COMPLETED
        material.save()
        return True

    except Exception as e:
        logger.error(f"Failed to index module material {material.id}: {e}")
        material.processing_status = ModuleMaterial.ProcessingStatus.FAILED
        material.save()
        return False
